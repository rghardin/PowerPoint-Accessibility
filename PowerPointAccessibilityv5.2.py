# -*- coding: utf-8 -*-
"""
PowerpointAccessibility.py
Main Streamlit application for PowerPoint accessibility improvements.
Uses LibreOffice for PDF generation (cross-platform).

@author: robert.hardin
"""

import requests
import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
import zipfile
import os
import tempfile
import shutil
import platform
import subprocess
from PIL import Image
import fitz

def call_models_api(api_key):
    """Fetch available models from TAMU AI API"""
    url = "https://chat-api.tamu.ai/openai/models"
    headers = {
        "accept": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    model_info = response.json()['data']
    id_list = []
    name_list = []
    for model in model_info:
        id_list.append(model['id'])
        name_list.append(model['name'])
    model_dict = dict(zip(name_list, id_list))
    return model_dict


def interact_with_model(api_key, chosen_model, my_query):
    """Send text-only query to LLM and get response"""
    url = "https://chat-api.tamu.ai/openai/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": chosen_model,
        "messages": [{"role": "user", "content": my_query}],
        "stream": False
    }
    response = requests.post(url, headers=headers, json=payload)
    return response.json()


def interact_with_model_vision(api_key, chosen_model, text_prompt, image_base64, image_format="png"):
    """
    Send query with image to vision-capable LLM and get response.
    Uses OpenAI vision API format with base64-encoded images.
    """
    url = "https://chat-api.tamu.ai/openai/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    
    mime_type_map = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp"
    }
    mime_type = mime_type_map.get(image_format.lower(), "image/png")
    
    message_content = [
        {"type": "text", "text": text_prompt},
        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_base64}"}}
    ]
    
    payload = {
        "model": chosen_model,
        "messages": [{"role": "user", "content": message_content}],
        "stream": False,
        "max_tokens": 300
    }
    
    response = requests.post(url, headers=headers, json=payload)
    return response.json()


def extract_llm_response(response_json):
    """Extract text content from LLM response"""
    try:
        return response_json['choices'][0]['message']['content']
    except (KeyError, IndexError):
        return ""


def check_vision_error(response_json):
    """Check if the response indicates a vision-related error."""
    try:
        if 'error' in response_json:
            error_msg = response_json['error'].get('message', str(response_json['error']))
            return True, error_msg
        return False, ""
    except:
        return False, ""


def image_to_base64(image_blob):
    """Convert image blob to base64 string"""
    return base64.b64encode(image_blob).decode('utf-8')


def get_image_format(image_blob):
    """Detect image format from blob header bytes"""
    if image_blob[:8] == b'\x89PNG\r\n\x1a\n':
        return "png"
    elif image_blob[:2] == b'\xff\xd8':
        return "jpeg"
    elif image_blob[:6] in (b'GIF87a', b'GIF89a'):
        return "gif"
    elif image_blob[:4] == b'RIFF' and image_blob[8:12] == b'WEBP':
        return "webp"
    elif image_blob[:4] == b'\x01\x00\x00\x00' or image_blob[:4] == b'\xd7\xcd\xc6\x9a':
        return "emf"
    else:
        return "png"


def convert_image_to_supported_format(image_blob):
    """
    Convert image blob to a supported format (PNG) if necessary.
    Uses LibreOffice for EMF/WMF conversion (cloud-compatible).
    Returns (converted_blob, format) tuple.
    """
    original_format = get_image_format(image_blob)
    supported_formats = ['png', 'jpeg', 'gif', 'webp']
    
    if original_format.lower() in supported_formats:
        return image_blob, original_format
    
    # Use LibreOffice for EMF/WMF conversion (replaces Wand/ImageMagick)
    if original_format.lower() in ['emf', 'wmf']:
        converted_blob = convert_emf_wmf_with_libreoffice(image_blob, original_format)
        if converted_blob is not None:
            return converted_blob, 'png'
        else:
            st.warning(f"LibreOffice conversion failed for {original_format.upper()} image")
    
    # Fallback to PIL for other formats
    try:
        img = Image.open(BytesIO(image_blob))
        
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGBA')
        elif img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')
        
        output_buffer = BytesIO()
        img.save(output_buffer, format='PNG')
        output_buffer.seek(0)
        
        return output_buffer.read(), 'png'
    
    except Exception as e:
        st.warning(f"Image conversion failed: {str(e)}")
        return image_blob, original_format


def convert_emf_wmf_with_libreoffice(image_blob, image_format):
    """
    Convert EMF/WMF image to PNG using LibreOffice.
    This is cloud-compatible and doesn't require ImageMagick.
    
    Args:
        image_blob: Binary image data
        image_format: 'emf' or 'wmf'
    
    Returns:
        PNG image blob or None if conversion failed
    """
    soffice = find_libreoffice()
    
    if not soffice:
        st.warning("LibreOffice not found for image conversion")
        return None
    
    temp_dir = None
    try:
        # Create temporary directory for conversion
        temp_dir = tempfile.mkdtemp()
        
        # Determine file extension
        ext = image_format.lower()
        if ext not in ['emf', 'wmf']:
            ext = 'emf'  # Default to EMF
        
        # Save image blob to temporary file
        input_path = os.path.join(temp_dir, f"temp_image.{ext}")
        with open(input_path, 'wb') as f:
            f.write(image_blob)
        
        # Use LibreOffice Draw to convert to PNG
        # LibreOffice can open EMF/WMF and export to various formats
        cmd = [
            soffice,
            '--headless',
            '--convert-to', 'png',
            '--outdir', temp_dir,
            input_path
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=30  # 30 second timeout for single image
        )
        
        # Look for the output PNG file
        output_path = os.path.join(temp_dir, "temp_image.png")
        
        if os.path.exists(output_path):
            with open(output_path, 'rb') as f:
                png_blob = f.read()
            return png_blob
        else:
            # LibreOffice might not directly convert EMF/WMF
            # Try alternative approach using a minimal ODP wrapper
            return convert_emf_wmf_via_odp(image_blob, image_format, temp_dir, soffice)
            
    except subprocess.TimeoutExpired:
        st.warning("LibreOffice image conversion timed out")
        return None
    except Exception as e:
        st.warning(f"LibreOffice image conversion error: {str(e)}")
        return None
    finally:
        # Clean up temporary directory
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except:
                pass


def convert_emf_wmf_via_odp(image_blob, image_format, temp_dir, soffice):
    """
    Convert EMF/WMF to PNG by embedding in an ODP document and exporting.
    This is a more reliable method for LibreOffice conversion.
    
    Args:
        image_blob: Binary image data
        image_format: 'emf' or 'wmf'
        temp_dir: Temporary directory for working files
        soffice: Path to LibreOffice executable
    
    Returns:
        PNG image blob or None if conversion failed
    """
    try:
        ext = image_format.lower()
        
        # Create a minimal ODP (OpenDocument Presentation) file with the image
        odp_path = os.path.join(temp_dir, "temp_presentation.odp")
        
        # Create ODP structure
        odp_content_dir = os.path.join(temp_dir, "odp_content")
        os.makedirs(odp_content_dir, exist_ok=True)
        os.makedirs(os.path.join(odp_content_dir, "Pictures"), exist_ok=True)
        
        # Save the image
        image_filename = f"image.{ext}"
        image_path = os.path.join(odp_content_dir, "Pictures", image_filename)
        with open(image_path, 'wb') as f:
            f.write(image_blob)
        
        # Determine MIME type
        mime_type = "image/x-emf" if ext == "emf" else "image/x-wmf"
        
        # Create manifest.xml
        manifest_content = f'''<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">
    <manifest:file-entry manifest:media-type="application/vnd.oasis.opendocument.presentation" manifest:full-path="/"/>
    <manifest:file-entry manifest:media-type="{mime_type}" manifest:full-path="Pictures/{image_filename}"/>
    <manifest:file-entry manifest:media-type="text/xml" manifest:full-path="content.xml"/>
</manifest:manifest>'''
        
        meta_inf_dir = os.path.join(odp_content_dir, "META-INF")
        os.makedirs(meta_inf_dir, exist_ok=True)
        with open(os.path.join(meta_inf_dir, "manifest.xml"), 'w', encoding='utf-8') as f:
            f.write(manifest_content)
        
        # Create content.xml with the image
        content_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<office:document-content 
    xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
    xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    xmlns:xlink="http://www.w3.org/1999/xlink"
    xmlns:svg="urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0"
    xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0">
    <office:body>
        <office:presentation>
            <draw:page draw:name="page1" draw:style-name="dp1" presentation:presentation-page-layout-name="AL1T0">
                <draw:frame draw:style-name="gr1" draw:layer="layout" svg:width="20cm" svg:height="15cm" svg:x="2cm" svg:y="2cm">
                    <draw:image xlink:href="Pictures/{image_filename}" xlink:type="simple" xlink:show="embed" xlink:actuate="onLoad"/>
                </draw:frame>
            </draw:page>
        </office:presentation>
    </office:body>
</office:document-content>'''
        
        with open(os.path.join(odp_content_dir, "content.xml"), 'w', encoding='utf-8') as f:
            f.write(content_xml)
        
        # Create the ODP file (ZIP archive)
        with zipfile.ZipFile(odp_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            # Add mimetype first (uncompressed)
            zf.writestr("mimetype", "application/vnd.oasis.opendocument.presentation")
            
            # Add other files
            for root, dirs, files in os.walk(odp_content_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arc_name = os.path.relpath(file_path, odp_content_dir)
                    zf.write(file_path, arc_name)
        
        # Convert ODP to PNG using LibreOffice
        cmd = [
            soffice,
            '--headless',
            '--convert-to', 'png',
            '--outdir', temp_dir,
            odp_path
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=30
        )
        
        # Look for the output PNG
        output_path = os.path.join(temp_dir, "temp_presentation.png")
        
        if os.path.exists(output_path):
            with open(output_path, 'rb') as f:
                png_blob = f.read()
            return png_blob
        
        return None
        
    except Exception as e:
        st.warning(f"ODP conversion method failed: {str(e)}")
        return None


def generate_image_alt_text(api_key, model_id, image_blob, context="", use_vision=True):
    """Generate alt text for an image using the LLM."""
    base_prompt = """Generate concise, descriptive alt text for this image in a PowerPoint presentation.
    The alt text should:
    - Be concise (under 125 characters if possible)
    - Describe the content and function of the image
    - Not start with "Image of" or "Picture of"
    - Be meaningful for screen reader users
    """
    
    context_addition = f"\n\nAdditional context from the slide: {context}" if context else ""
    
    if use_vision and image_blob is not None:
        try:
            converted_blob, image_format = convert_image_to_supported_format(image_blob)
            
            original_format = get_image_format(image_blob)
            if original_format != image_format:
                st.info(f"Converted image from {original_format.upper()} to {image_format.upper()}")
            
            image_base64 = image_to_base64(converted_blob)
            vision_prompt = base_prompt + context_addition + "\n\nPlease analyze the image and provide only the alt text, nothing else."
            
            response = interact_with_model_vision(api_key, model_id, vision_prompt, image_base64, image_format)
            
            is_error, error_msg = check_vision_error(response)
            if not is_error:
                alt_text = extract_llm_response(response).strip()
                if alt_text:
                    return alt_text, True
            else:
                st.warning(f"Vision not supported: {error_msg}. Using text-based generation.")
        except Exception as e:
            st.warning(f"Vision request failed: {str(e)}. Using text-based generation.")
    
    fallback_prompt = f"""{base_prompt}
    Since I cannot see the image directly, please generate appropriate alt text based on the context from the slide.
    Context from the slide: {context if context else 'No additional context available'}
    Please provide only the alt text, nothing else."""
    
    response = interact_with_model(api_key, model_id, fallback_prompt)
    return extract_llm_response(response).strip(), False


def generate_table_alt_text(api_key, model_id, table_data, context=""):
    """Generate alt text/summary for a table using the LLM"""
    prompt = f"""Generate a concise summary description for a table in a PowerPoint presentation.
    Table content: {table_data}
    Context from the slide: {context if context else 'No additional context available'}
    The summary should be under 200 characters. Please provide only the summary, nothing else."""
    
    response = interact_with_model(api_key, model_id, prompt)
    return extract_llm_response(response).strip()


def extract_table_data(table):
    """Extract text content from a table shape"""
    table_text = []
    for row in table.rows:
        row_data = [cell.text.strip() for cell in row.cells]
        table_text.append(" | ".join(row_data))
    return "\n".join(table_text)


def get_slide_context(slide):
    """Extract text content from a slide for context"""
    text_content = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_content.append(text)
    return " ".join(text_content[:500])


def set_image_alt_text(shape, alt_text, title=""):
    """Set alt text for an image shape"""
    nvSpPr = shape._element._nvXxPr
    cNvPr = nvSpPr.cNvPr
    cNvPr.set('descr', alt_text)
    if title:
        cNvPr.set('title', title)


def set_table_alt_text(shape, alt_text, title=""):
    """Set alt text for a table shape"""
    try:
        nvSpPr = shape._element._nvXxPr
        cNvPr = nvSpPr.cNvPr
        cNvPr.set('descr', alt_text)
        if title:
            cNvPr.set('title', title)
    except Exception:
        pass


def ensure_slide_titles(prs):
    """Count slides without titles"""
    untitled_count = 0
    for slide in prs.slides:
        has_title = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                if shape.placeholder_format.type == 1 and shape.text.strip():
                    has_title = True
                    break
        if not has_title:
            untitled_count += 1
    return untitled_count


def set_document_properties(prs, filename):
    """Set document properties for accessibility"""
    title = os.path.splitext(filename)[0]
    title = title.replace('_', ' ').replace('-', ' ').title()
    
    if not prs.core_properties.title:
        prs.core_properties.title = title
    return title


def process_powerpoint(uploaded_file, api_key, model_id, use_vision=True):
    """Process a PowerPoint file for accessibility"""
    prs = Presentation(BytesIO(uploaded_file.read()))
    uploaded_file.seek(0)
    
    log_messages = []
    alt_text_log = []
    images_processed = 0
    tables_processed = 0
    vision_used_count = 0
    text_fallback_count = 0
    
    doc_title = set_document_properties(prs, uploaded_file.name)
    log_messages.append(f"✓ Set document title: '{doc_title}'")
    
    untitled = ensure_slide_titles(prs)
    if untitled > 0:
        log_messages.append(f"⚠ Warning: {untitled} slide(s) without titles detected")
    
    total_slides = len(prs.slides)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_context = get_slide_context(slide)
        
        for shape in slide.shapes:
            has_image = False
            image_blob = None
            try:
                if hasattr(shape, 'image') and shape.image is not None:
                    has_image = True
                    image_blob = shape.image.blob
            except:
                pass
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or has_image:
                try:
                    existing_alt = ""
                    try:
                        existing_alt = shape._element._nvXxPr.cNvPr.get('descr', '')
                    except:
                        pass
                    
                    if not existing_alt or existing_alt.strip() == "":
                        alt_text, used_vision = generate_image_alt_text(
                            api_key, model_id, image_blob, slide_context, use_vision=use_vision
                        )
                        
                        if used_vision:
                            vision_used_count += 1
                            method = "🖼️ Vision"
                        else:
                            text_fallback_count += 1
                            method = "📝 Text-only"
                        
                        st.info(f"**Slide {slide_idx + 1} - {method}:**\n{alt_text}")
                        alt_text_log.append({
                            'slide': slide_idx + 1, 'type': 'Image',
                            'method': 'Vision' if used_vision else 'Text-only', 'alt_text': alt_text
                        })
                        
                        set_image_alt_text(shape, alt_text)
                        images_processed += 1
                        log_messages.append(f"  Slide {slide_idx + 1}: Added alt text to image ({method})")
                except Exception as e:
                    log_messages.append(f"  Slide {slide_idx + 1}: Error processing image - {str(e)}")
            
            elif shape.has_table:
                try:
                    table_data = extract_table_data(shape.table)
                    alt_text = generate_table_alt_text(api_key, model_id, table_data, slide_context)
                    
                    st.info(f"**Slide {slide_idx + 1} - Table:**\n{alt_text}")
                    alt_text_log.append({
                        'slide': slide_idx + 1, 'type': 'Table',
                        'method': 'Text analysis', 'alt_text': alt_text
                    })
                    
                    set_table_alt_text(shape, alt_text)
                    tables_processed += 1
                    log_messages.append(f"  Slide {slide_idx + 1}: Added description to table")
                except Exception as e:
                    log_messages.append(f"  Slide {slide_idx + 1}: Error processing table - {str(e)}")
    
    log_messages.append(f"✓ Processed {images_processed} images ({vision_used_count} with vision, {text_fallback_count} text-only) and {tables_processed} tables")
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output, log_messages, alt_text_log

def find_libreoffice():
    """Find LibreOffice executable path"""
    if platform.system() == "Windows":
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None
    else:
        # Linux/Mac - check if soffice is in PATH
        try:
            result = subprocess.run(['which', 'soffice'], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout.strip()
        except:
            pass
        
        # Check common Linux paths
        possible_paths = [
            '/usr/bin/soffice',
            '/usr/bin/libreoffice',
            '/usr/local/bin/soffice',
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None


def convert_pptx_to_slides_pdf_libreoffice(pptx_path, output_dir):
    """
    Convert PowerPoint to slides PDF using LibreOffice.
    Returns the path to the generated PDF or None if failed.
    """
    soffice = find_libreoffice()
    
    if not soffice:
        st.error("LibreOffice not found. Please install LibreOffice to enable PDF conversion.")
        return None
    
    try:
        cmd = [
            soffice,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            pptx_path
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120
        )
        
        if result.returncode != 0:
            st.warning(f"LibreOffice conversion warning: {result.stderr}")
        
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
        
        if os.path.exists(pdf_path):
            return pdf_path
        else:
            st.warning(f"PDF file not found after conversion: {pdf_path}")
            return None
            
    except subprocess.TimeoutExpired:
        st.warning("LibreOffice conversion timed out")
        return None
    except Exception as e:
        st.warning(f"LibreOffice conversion failed: {str(e)}")
        return None


def convert_pptx_to_handout_pdf_libreoffice(pptx_path, output_dir, slides_per_page=3):
    """
    Convert PowerPoint to handout PDF using LibreOffice with filter options.
    
    Args:
        pptx_path: Path to the input PowerPoint file
        output_dir: Directory to save the output PDF
        slides_per_page: Number of slides per page (2, 3, 4, or 6)
    
    Returns:
        Path to the generated PDF or None if failed
    """
    soffice = find_libreoffice()
    
    if not soffice:
        st.error("LibreOffice not found. Please install LibreOffice to enable PDF conversion.")
        return None
    
    # Map slides per page to LibreOffice handout type
    # LibreOffice SlidesPerPage values: 1, 2, 3, 4, 6, 9
    handout_map = {2: 2, 3: 3, 4: 4, 6: 6}
    slides_per_page_value = handout_map.get(slides_per_page, 3)
    
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    
    try:
        # LibreOffice PDF export filter options for handouts
        # PageContentType: 0=slides, 1=notes, 2=handouts, 3=outline
        # SlidesPerPage: number of slides per handout page
        filter_options = (
            f"impress_pdf_Export:"
            f"PageContentType=2,"  # Handouts
            f"SlidesPerPage={slides_per_page_value},"
            f"IsAddStream=false"
        )
        
        cmd = [
            soffice,
            '--headless',
            '--convert-to', f'pdf:{filter_options}',
            '--outdir', output_dir,
            pptx_path
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120
        )
        
        if result.returncode != 0:
            st.warning(f"LibreOffice handout conversion warning: {result.stderr}")
        
        # LibreOffice creates the PDF with the same base name
        pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
        
        if os.path.exists(pdf_path):
            # Rename to indicate it's a handout
            handout_pdf_path = os.path.join(output_dir, f"{base_name}_handouts.pdf")
            os.rename(pdf_path, handout_pdf_path)
            return handout_pdf_path
        else:
            st.warning(f"Handout PDF file not found after conversion")
            return None
            
    except subprocess.TimeoutExpired:
        st.warning("LibreOffice handout conversion timed out")
        return None
    except Exception as e:
        st.warning(f"LibreOffice handout conversion failed: {str(e)}")
        return None

def create_handout_from_slides_pdf(slides_pdf_bytes, slides_per_page=3):
    """
    Create handout-style PDF by arranging multiple slides per page.
    Uses PyMuPDF to rearrange slides from an existing PDF.
    
    For 3 slides per page: Creates PowerPoint-style handouts with slides
    on the left and horizontal lines for notes on the right.
    
    Args:
        slides_pdf_bytes: BytesIO containing the slides PDF
        slides_per_page: Number of slides per page (2, 3, 4, or 6)
    
    Returns:
        BytesIO containing the handout PDF
    """
    slides_pdf_bytes.seek(0)
    doc = fitz.open(stream=slides_pdf_bytes.read(), filetype="pdf")
    new_doc = fitz.open()
    
    # Page size (Letter: 612 x 792 points)
    page_width = 612
    page_height = 792
    
    # Margins
    margin_left = 36      # 0.5 inch
    margin_right = 36     # 0.5 inch
    margin_top = 54       # 0.75 inch
    margin_bottom = 54    # 0.75 inch
    
    # Define layouts for different slides per page
    # Each tuple is (x1, y1, x2, y2) as fractions of page dimensions
    layouts = {
        2: [
            (0.05, 0.05, 0.95, 0.47),   # Top slide
            (0.05, 0.53, 0.95, 0.95)    # Bottom slide
        ],
        3: [
            # PowerPoint-style: slides on LEFT side
            (0.05, 0.05, 0.45, 0.30),   # Top left
            (0.05, 0.37, 0.45, 0.62),   # Middle left
            (0.05, 0.69, 0.45, 0.94)    # Bottom left
        ],
        4: [
            (0.03, 0.03, 0.48, 0.47),   # Top left
            (0.52, 0.03, 0.97, 0.47),   # Top right
            (0.03, 0.53, 0.48, 0.97),   # Bottom left
            (0.52, 0.53, 0.97, 0.97)    # Bottom right
        ],
        6: [
            (0.03, 0.03, 0.48, 0.30),   # Row 1 left
            (0.52, 0.03, 0.97, 0.30),   # Row 1 right
            (0.03, 0.35, 0.48, 0.62),   # Row 2 left
            (0.52, 0.35, 0.97, 0.62),   # Row 2 right
            (0.03, 0.68, 0.48, 0.95),   # Row 3 left
            (0.52, 0.68, 0.97, 0.95)    # Row 3 right
        ]
    }
    
    # Note line areas for 3-slides-per-page layout (right side of each slide)
    # Each tuple is (x_start, y_start, x_end, y_end) for the notes area
    notes_areas_3 = [
        (0.50, 0.05, 0.95, 0.30),   # Notes area for top slide
        (0.50, 0.37, 0.95, 0.62),   # Notes area for middle slide
        (0.50, 0.69, 0.95, 0.94)    # Notes area for bottom slide
    ]
    
    positions = layouts.get(slides_per_page, layouts[3])
    num_positions = len(positions)
    
    # Process slides in groups
    for group_start in range(0, len(doc), num_positions):
        # Create new page
        new_page = new_doc.new_page(width=page_width, height=page_height)
        
        # Add slides to this page
        for i, pos in enumerate(positions):
            slide_idx = group_start + i
            if slide_idx >= len(doc):
                break
            
            # Calculate destination rectangle for the slide
            rect = fitz.Rect(
                pos[0] * page_width,
                pos[1] * page_height,
                pos[2] * page_width,
                pos[3] * page_height
            )
            
            # Insert the slide page into the rectangle
            new_page.show_pdf_page(rect, doc, slide_idx)
            
            # Add border around the slide
            new_page.draw_rect(rect, color=(0.5, 0.5, 0.5), width=0.75)
            
            # For 3 slides per page, add note lines on the right side
            if slides_per_page == 3:
                draw_note_lines(new_page, notes_areas_3[i], page_width, page_height)
    
    # Handle case where we have remaining slides that don't fill the last page
    # (already handled by the break in the loop above)
    
    # Save to BytesIO
    output = BytesIO()
    new_doc.save(output)
    new_doc.close()
    doc.close()
    
    output.seek(0)
    return output


def draw_note_lines(page, notes_area, page_width, page_height):
    """
    Draw horizontal lines for note-taking in the specified area.
    
    Args:
        page: PyMuPDF page object
        notes_area: Tuple (x1, y1, x2, y2) as fractions of page dimensions
        page_width: Page width in points
        page_height: Page height in points
    """
    # Calculate the notes area rectangle
    x_start = notes_area[0] * page_width
    y_start = notes_area[1] * page_height
    x_end = notes_area[2] * page_width
    y_end = notes_area[3] * page_height
    
    # Line spacing (approximately 18 points = 0.25 inch, typical ruled paper)
    line_spacing = 18
    
    # Add a small top margin before starting lines
    top_padding = 10
    
    # Line color (light gray)
    line_color = (0.75, 0.75, 0.75)
    
    # Calculate number of lines that fit in the area
    available_height = y_end - y_start - top_padding
    num_lines = int(available_height / line_spacing)
    
    # Draw horizontal lines
    current_y = y_start + top_padding
    for _ in range(num_lines):
        # Draw a horizontal line
        page.draw_line(
            fitz.Point(x_start, current_y),
            fitz.Point(x_end, current_y),
            color=line_color,
            width=0.5
        )
        current_y += line_spacing

def convert_pptx_to_pdf_batch_libreoffice(pptx_files_dict, temp_dir, generate_slides=False, generate_handouts=False, slides_per_page=3):
    """
    Convert multiple PowerPoint files to PDF using LibreOffice.
    Creates handouts by post-processing with PyMuPDF.
    """
    slides_pdf_dict = {}
    handouts_pdf_dict = {}
    
    if not generate_slides and not generate_handouts:
        return slides_pdf_dict, handouts_pdf_dict
    
    soffice = find_libreoffice()
    if not soffice:
        st.error("LibreOffice not found. PDF generation is disabled.")
        return slides_pdf_dict, handouts_pdf_dict
    
    for filename, pptx_bytes in pptx_files_dict.items():
        base_name = os.path.splitext(filename)[0]
        
        # Save PPTX to temp file
        temp_pptx_path = os.path.join(temp_dir, filename)
        pptx_bytes.seek(0)
        with open(temp_pptx_path, 'wb') as f:
            f.write(pptx_bytes.read())
        pptx_bytes.seek(0)
        
        # Convert to slides PDF using LibreOffice
        try:
            cmd = [
                soffice,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                temp_pptx_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            # Find the generated PDF
            slides_pdf_path = os.path.join(temp_dir, f"{base_name}.pdf")
            
            if os.path.exists(slides_pdf_path):
                # Read slides PDF
                with open(slides_pdf_path, 'rb') as f:
                    slides_pdf_bytes = BytesIO(f.read())
                
                # Store slides PDF if requested
                if generate_slides:
                    slides_pdf_bytes.seek(0)
                    slides_pdf_dict[f"{base_name}_slides.pdf"] = BytesIO(slides_pdf_bytes.read())
                    st.write(f"  ✓ Slides PDF created for {filename}")
                
                # Create handouts PDF if requested using PyMuPDF
                if generate_handouts:
                    slides_pdf_bytes.seek(0)
                    handouts_pdf = create_handout_from_slides_pdf(slides_pdf_bytes, slides_per_page)
                    handouts_pdf_dict[f"{base_name}_handouts.pdf"] = handouts_pdf
                    st.write(f"  ✓ Handouts PDF ({slides_per_page} slides/page) created for {filename}")
            else:
                st.warning(f"  ✗ Failed to create PDF for {filename}")
                
        except subprocess.TimeoutExpired:
            st.warning(f"PDF conversion timed out for {filename}")
        except Exception as e:
            st.warning(f"Error converting {filename} to PDF: {str(e)}")
    
    return slides_pdf_dict, handouts_pdf_dict

def create_zip_file(files_dict):
    """Create a ZIP file from a dictionary of {filename: bytes_io}"""
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for filename, file_bytes in files_dict.items():
            file_bytes.seek(0)
            zip_file.writestr(filename, file_bytes.read())
    zip_buffer.seek(0)
    return zip_buffer


# Main Streamlit App
st.title("PowerPoint Digital Accessibility Improvements and File Save Utility")

st.markdown("""
This application improves the digital accessibility of PowerPoint files by generating alt text
for images and tables and enables batch saving of PDF slides.
""")

api_key = st.text_input("TAMU API Key", value=None, type="password")

if api_key is not None and api_key != "":
    try:
        model_dict = call_models_api(api_key)
        selected_model_name = st.selectbox(
            "Pick a large language model to use for generating alt text", 
            list(model_dict), key=25
        )
        selected_model_id = model_dict[selected_model_name]
        
        # Selection for completed files
        st.subheader("Output Options")
        
        # Check if LibreOffice is available (works on both Windows and Linux)
        libreoffice_available = find_libreoffice() is not None
        
        if libreoffice_available:
            slides = st.checkbox("Select to output slides in PDF format")
            handouts = st.checkbox("Select to output handouts in PDF format")
            
            if handouts:
                slides_per_page = st.radio("Select number of slides per page", [2, 3, 4, 6], index=1)
            else:
                slides_per_page = 3
        else:
            st.warning("LibreOffice not found. PDF generation is disabled. Install LibreOffice to enable this feature.")
            slides = False
            handouts = False
            slides_per_page = 3
                
        # File uploader
        uploaded_files = st.file_uploader(
            "Choose PowerPoint files to modify", 
            type="pptx", 
            accept_multiple_files=True
        )

        # Process button
        if uploaded_files:
            st.write(f"**{len(uploaded_files)} file(s) uploaded**")
            
            if st.button("Process PowerPoint Files", type="primary"):
                processed_pptx = {}
                processed_slides_pdf = {}
                processed_handouts_pdf = {}
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # Step 1: Process all PowerPoint files for accessibility
                with st.expander("Processing Details - Accessibility", expanded=True):
                    for i, uploaded_file in enumerate(uploaded_files):
                        status_text.text(f"Processing accessibility: {uploaded_file.name}")
                        st.write(f"**Processing: {uploaded_file.name}**")
                        
                        try:
                            # Process PowerPoint for accessibility
                            modified_pptx, log_messages, alt_text_log = process_powerpoint(
                                uploaded_file, api_key, selected_model_id, use_vision=True
                            )
                            
                            # Display log messages
                            for msg in log_messages:
                                st.write(msg)
                            
                            # Store processed PPTX
                            base_name = os.path.splitext(uploaded_file.name)[0]
                            processed_pptx[f"{base_name}_accessible.pptx"] = modified_pptx
                            
                            st.success(f"✓ Completed accessibility processing: {uploaded_file.name}")
                            
                        except Exception as e:
                            st.error(f"Error processing {uploaded_file.name}: {str(e)}")
                        
                        progress_bar.progress((i + 1) / len(uploaded_files) * 0.5)
                
                # Step 2: Generate PDFs if requested using LibreOffice
                if slides or handouts:
                    with st.expander("Processing Details - PDF Generation", expanded=True):
                        status_text.text("Generating PDFs...")
                        st.write("**Generating PDF files using LibreOffice...**")
                        
                        # Create temp directory for PDF generation
                        temp_dir = tempfile.mkdtemp()
                        
                        try:
                            processed_slides_pdf, processed_handouts_pdf = convert_pptx_to_pdf_batch_libreoffice(
                                processed_pptx,
                                temp_dir,
                                generate_slides=slides,
                                generate_handouts=handouts,
                                slides_per_page=slides_per_page
                                )
            
                            if slides:
                                st.write(f"✓ Generated {len(processed_slides_pdf)} slides PDF(s)")
                            if handouts:
                                st.write(f"✓ Generated {len(processed_handouts_pdf)} handouts PDF(s)")
                
                        except Exception as e:
                            st.error(f"Error generating PDFs: {str(e)}")
                
                        finally:
                            # Clean up temp files
                            try:
                                shutil.rmtree(temp_dir)
                            except:
                                pass
                        
                        progress_bar.progress(1.0)
                
                status_text.text("Processing complete!")
                progress_bar.progress(1.0)
                
                # Store in session state to persist after button click
                st.session_state['processed_pptx'] = processed_pptx
                st.session_state['processed_slides_pdf'] = processed_slides_pdf
                st.session_state['processed_handouts_pdf'] = processed_handouts_pdf
                st.session_state['processing_complete'] = True
        
        # Download section - check if processing is complete
        if st.session_state.get('processing_complete', False):
            st.subheader("📥 Download Processed Files")
            
            processed_pptx = st.session_state.get('processed_pptx', {})
            processed_slides_pdf = st.session_state.get('processed_slides_pdf', {})
            processed_handouts_pdf = st.session_state.get('processed_handouts_pdf', {})
            
            if len(processed_pptx) == 1:
                # Single file - provide individual downloads
                st.write("**Download individual files:**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    for filename, file_bytes in processed_pptx.items():
                        file_bytes.seek(0)
                        st.download_button(
                            label="📄 Download Accessible PPTX",
                            data=file_bytes,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="download_pptx"
                        )
                
                with col2:
                    if processed_slides_pdf:
                        for filename, file_bytes in processed_slides_pdf.items():
                            file_bytes.seek(0)
                            st.download_button(
                                label="📑 Download Slides PDF",
                                data=file_bytes,
                                file_name=filename,
                                mime="application/pdf",
                                key="download_slides"
                            )
                    else:
                        st.write("*No slides PDF*")
                
                with col3:
                    if processed_handouts_pdf:
                        for filename, file_bytes in processed_handouts_pdf.items():
                            file_bytes.seek(0)
                            st.download_button(
                                label="📋 Download Handouts PDF",
                                data=file_bytes,
                                file_name=filename,
                                mime="application/pdf",
                                key="download_handouts"
                            )
                    else:
                        st.write("*No handouts PDF*")
                        
            elif len(processed_pptx) > 1:
                # Multiple files - provide ZIP downloads
                st.write("**Download ZIP archives:**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if processed_pptx:
                        pptx_zip = create_zip_file(processed_pptx)
                        st.download_button(
                            label="📦 Download All PPTX (ZIP)",
                            data=pptx_zip,
                            file_name="accessible_presentations.zip",
                            mime="application/zip",
                            key="download_pptx_zip"
                        )
                
                with col2:
                    if processed_slides_pdf:
                        slides_zip = create_zip_file(processed_slides_pdf)
                        st.download_button(
                            label="📦 Download All Slides PDF (ZIP)",
                            data=slides_zip,
                            file_name="slides_pdfs.zip",
                            mime="application/zip",
                            key="download_slides_zip"
                        )
                    else:
                        st.write("*No slides PDFs*")
               
                with col3:
                    if processed_handouts_pdf:
                        handouts_zip = create_zip_file(processed_handouts_pdf)
                        st.download_button(
                            label="📦 Download All Handouts PDF (ZIP)",
                            data=handouts_zip,
                            file_name="handouts_pdfs.zip",
                            mime="application/zip",
                            key="download_handouts_zip"
                        )
                    else:
                        st.write("*No handouts PDFs*")
                        
                # Option to download everything in one ZIP
                st.write("---")
                all_files = {**processed_pptx, **processed_slides_pdf, **processed_handouts_pdf}
                if all_files:
                    all_zip = create_zip_file(all_files)
                    st.download_button(
                        label="📦 Download All Files (Single ZIP)",
                        data=all_zip,
                        file_name="all_processed_files.zip",
                        mime="application/zip",
                        key="download_all_zip"
                    )
            
            # Reset button
            if st.button("🔄 Process New Files"):
                st.session_state['processing_complete'] = False
                st.session_state['processed_pptx'] = {}
                st.session_state['processed_slides_pdf'] = {}
                st.session_state['processed_handouts_pdf'] = {}
                st.rerun()
                
    except Exception as e:
        st.error(f"Error connecting to API: {str(e)}")
        st.write("Please check your API key and try again.")

else:
    st.warning("Please enter your TAMU API Key to continue.")
    st.markdown("""
    ### Instructions:
    1. Enter your TAMU AI API key above
    2. Select a language model for use in generating alt text
    3. Select to output slides as PDF (requires LibreOffice)
    4. Upload PowerPoint files to process
    5. Click 'Process PowerPoint Files' to make them accessible
    6. Download the modified files individually or as a ZIP
    """)