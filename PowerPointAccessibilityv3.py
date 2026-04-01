# -*- coding: utf-8 -*-
"""
PowerpointAccessibility.py
Main Streamlit application for PowerPoint accessibility improvements.
Uses subprocess to call pdf_converter.py for PDF generation on Windows.

@author: robert.hardin
"""

import requests
import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
import base64
import zipfile
import os
import tempfile
import shutil
import platform
import subprocess
import json
from PIL import Image


def call_models_api(api_key):
    """Fetch available models from TAMU AI API"""
    url = "https://chat-api.tamu.ai/openai/models"
    headers = {
        "accept": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    model_info = response.json()['data']
    id_list = []
    name_list = []
    for model in model_info:
        id_list.append(model['id'])
        name_list.append(model['name'])
    model_dict = dict(zip(name_list, id_list))
    return model_dict


def interact_with_model(api_key, chosen_model, my_query):
    """Send text-only query to LLM and get response"""
    url = "https://chat-api.tamu.ai/openai/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": chosen_model,
        "messages": [{"role": "user", "content": my_query}],
        "stream": False
    }
    response = requests.post(url, headers=headers, json=payload)
    return response.json()


def interact_with_model_vision(api_key, chosen_model, text_prompt, image_base64, image_format="png"):
    """
    Send query with image to vision-capable LLM and get response.
    Uses OpenAI vision API format with base64-encoded images.
    """
    url = "https://chat-api.tamu.ai/openai/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    
    # Determine MIME type based on image format
    mime_type_map = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp"
    }
    mime_type = mime_type_map.get(image_format.lower(), "image/png")
    
    # Build message content with text and image
    message_content = [
        {
            "type": "text",
            "text": text_prompt
        },
        {
            "type": "image_url",
            "image_url": {
                "url": f"data:{mime_type};base64,{image_base64}"
            }
        }
    ]
    
    payload = {
        "model": chosen_model,
        "messages": [{"role": "user", "content": message_content}],
        "stream": False,
        "max_tokens": 300
    }
    
    response = requests.post(url, headers=headers, json=payload)
    return response.json()


def extract_llm_response(response_json):
    """Extract text content from LLM response"""
    try:
        return response_json['choices'][0]['message']['content']
    except (KeyError, IndexError):
        return ""


def check_vision_error(response_json):
    """
    Check if the response indicates a vision-related error.
    Returns (is_error, error_message) tuple.
    """
    try:
        if 'error' in response_json:
            error_msg = response_json['error'].get('message', str(response_json['error']))
            vision_error_keywords = [
                'image', 'vision', 'multimodal', 'does not support',
                'invalid', 'content type', 'unsupported'
            ]
            if any(keyword in error_msg.lower() for keyword in vision_error_keywords):
                return True, error_msg
            return True, error_msg
        return False, ""
    except:
        return False, ""


def image_to_base64(image_blob):
    """Convert image blob to base64 string"""
    return base64.b64encode(image_blob).decode('utf-8')


def get_image_format(image_blob):
    """Detect image format from blob header bytes"""
    if image_blob[:8] == b'\x89PNG\r\n\x1a\n':
        return "png"
    elif image_blob[:2] == b'\xff\xd8':
        return "jpeg"
    elif image_blob[:6] in (b'GIF87a', b'GIF89a'):
        return "gif"
    elif image_blob[:4] == b'RIFF' and image_blob[8:12] == b'WEBP':
        return "webp"
    elif image_blob[:4] == b'\x01\x00\x00\x00' or image_blob[:4] == b'\xd7\xcd\xc6\x9a':
        return "emf"
    else:
        return "png"


def convert_image_to_supported_format(image_blob):
    """
    Convert image blob to a supported format (PNG) if necessary.
    Returns (converted_blob, format) tuple.
    """
    original_format = get_image_format(image_blob)
    supported_formats = ['png', 'jpeg', 'gif', 'webp']
    
    if original_format.lower() in supported_formats:
        return image_blob, original_format
    
    try:
        img = Image.open(BytesIO(image_blob))
        
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGBA')
        elif img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')
        
        output_buffer = BytesIO()
        img.save(output_buffer, format='PNG')
        output_buffer.seek(0)
        
        return output_buffer.read(), 'png'
    
    except Exception as e:
        return image_blob, original_format


def generate_image_alt_text(api_key, model_id, image_blob, context="", use_vision=True):
    """
    Generate alt text for an image using the LLM.
    Attempts to use vision capabilities first, falls back to text-only if not supported.
    
    Returns: (alt_text, used_vision) tuple
    """
    base_prompt = """Generate concise, descriptive alt text for this image in a PowerPoint presentation.
    The alt text should:
    - Be concise (under 125 characters if possible)
    - Describe the content and function of the image
    - Not start with "Image of" or "Picture of"
    - Be meaningful for screen reader users
    """
    
    context_addition = f"\n\nAdditional context from the slide: {context}" if context else ""
    
    if use_vision and image_blob is not None:
        try:
            converted_blob, image_format = convert_image_to_supported_format(image_blob)
            
            original_format = get_image_format(image_blob)
            if original_format != image_format:
                st.info(f"Converted image from {original_format.upper()} to {image_format.upper()}")
            
            image_base64 = image_to_base64(converted_blob)
            
            vision_prompt = base_prompt + context_addition + "\n\nPlease analyze the image and provide only the alt text, nothing else."
            
            response = interact_with_model_vision(api_key, model_id, vision_prompt, image_base64, image_format)
            
            is_error, error_msg = check_vision_error(response)
            if is_error:
                st.warning(f"Vision not supported by this model: {error_msg}. Falling back to text-based generation.")
            else:
                alt_text = extract_llm_response(response).strip()
                if alt_text:
                    return alt_text, True
        except Exception as e:
            st.warning(f"Vision request failed: {str(e)}. Falling back to text-based generation.")
    
    fallback_prompt = f"""{base_prompt}
    
    Since I cannot see the image directly, please generate appropriate alt text based on the context from the slide.
    
    Context from the slide: {context if context else 'No additional context available'}
    
    Please provide only the alt text, nothing else. If context is insufficient, provide generic but helpful alt text."""
    
    response = interact_with_model(api_key, model_id, fallback_prompt)
    alt_text = extract_llm_response(response).strip()
    
    return alt_text, False


def generate_table_alt_text(api_key, model_id, table_data, context=""):
    """Generate alt text/summary for a table using the LLM"""
    prompt = f"""Generate a concise summary description for a table in a PowerPoint presentation.
    
    Table content:
    {table_data}
    
    Context from the slide: {context if context else 'No additional context available'}
    
    The summary should:
    - Briefly describe what the table contains
    - Mention the number of rows and columns
    - Highlight key information
    - Be under 200 characters
    
    Please provide only the summary, nothing else."""
    
    response = interact_with_model(api_key, model_id, prompt)
    return extract_llm_response(response).strip()


def extract_table_data(table):
    """Extract text content from a table shape"""
    table_text = []
    for row_idx, row in enumerate(table.rows):
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        table_text.append(" | ".join(row_data))
    return "\n".join(table_text)


def get_slide_context(slide):
    """Extract text content from a slide for context"""
    text_content = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_content.append(text)
    return " ".join(text_content[:500])


def set_image_alt_text(shape, alt_text, title=""):
    """Set alt text for an image shape"""
    nvSpPr = shape._element._nvXxPr
    cNvPr = nvSpPr.cNvPr
    cNvPr.set('descr', alt_text)
    if title:
        cNvPr.set('title', title)


def set_table_alt_text(shape, alt_text, title=""):
    """Set alt text for a table shape"""
    try:
        nvSpPr = shape._element._nvXxPr
        cNvPr = nvSpPr.cNvPr
        cNvPr.set('descr', alt_text)
        if title:
            cNvPr.set('title', title)
    except Exception:
        pass


def ensure_slide_titles(prs):
    """Ensure all slides have titles for accessibility"""
    untitled_count = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        has_title = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                if shape.placeholder_format.type == 1:
                    if shape.text.strip():
                        has_title = True
                        break
        if not has_title:
            untitled_count += 1
    return untitled_count

def set_document_properties(prs, filename):
    """Set document properties for accessibility"""
    title = os.path.splitext(filename)[0]
    title = title.replace('_', ' ').replace('-', ' ').title()
    
    core_props = prs.core_properties
    if not core_props.title:
        core_props.title = title
    
    return title


def process_powerpoint(uploaded_file, api_key, model_id, use_vision=True, progress_callback=None):
    """Process a PowerPoint file for accessibility"""
    prs = Presentation(BytesIO(uploaded_file.read()))
    uploaded_file.seek(0)
    
    log_messages = []
    alt_text_log = []  # For troubleshooting output
    images_processed = 0
    tables_processed = 0
    vision_used_count = 0
    text_fallback_count = 0
    
    # Set document title
    doc_title = set_document_properties(prs, uploaded_file.name)
    log_messages.append(f"✓ Set document title: '{doc_title}'")
    
    # Check for untitled slides
    untitled = ensure_slide_titles(prs)
    if untitled > 0:
        log_messages.append(f"⚠ Warning: {untitled} slide(s) without titles detected")
    
    total_slides = len(prs.slides)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_context = get_slide_context(slide)
        
        for shape in slide.shapes:
            # Check if shape contains an image
            has_image = False
            image_blob = None
            try:
                if hasattr(shape, 'image') and shape.image is not None:
                    has_image = True
                    image_blob = shape.image.blob
            except:
                pass
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or has_image:
                try:
                    # Check if alt text already exists
                    existing_alt = ""
                    try:
                        existing_alt = shape._element._nvXxPr.cNvPr.get('descr', '')
                    except:
                        pass
                    
                    if not existing_alt or existing_alt.strip() == "":
                        # Generate alt text using LLM (with vision if available)
                        alt_text, used_vision = generate_image_alt_text(
                            api_key, model_id, 
                            image_blob,
                            slide_context,
                            use_vision=use_vision
                        )
                        
                        # Track vision usage
                        if used_vision:
                            vision_used_count += 1
                            method = "🖼️ Vision"
                        else:
                            text_fallback_count += 1
                            method = "📝 Text-only"
                                                   
                        # Display and log the alt text
                        st.info(f"**Slide {slide_idx + 1} - {method}:**\n{alt_text}")
                        alt_text_log.append({
                            'slide': slide_idx + 1,
                            'type': 'Image',
                            'method': 'Vision' if used_vision else 'Text-only',
                            'alt_text': alt_text
                        })
                        
                        set_image_alt_text(shape, alt_text)
                        images_processed += 1
                        log_messages.append(f"  Slide {slide_idx + 1}: Added alt text to image ({method})")
                except Exception as e:
                    log_messages.append(f"  Slide {slide_idx + 1}: Error processing image - {str(e)}")
            
            # Process tables
            elif shape.has_table:
                try:
                    table_data = extract_table_data(shape.table)
                    alt_text = generate_table_alt_text(api_key, model_id, table_data, slide_context)
                    
                    st.info(f"**Slide {slide_idx + 1} - Table:**\n{alt_text}")
                    alt_text_log.append({
                        'slide': slide_idx + 1,
                        'type': 'Table',
                        'method': 'Text analysis',
                        'alt_text': alt_text
                    })
                    
                    set_table_alt_text(shape, alt_text)
                    tables_processed += 1
                    log_messages.append(f"  Slide {slide_idx + 1}: Added description to table")
                except Exception as e:
                    log_messages.append(f"  Slide {slide_idx + 1}: Error processing table - {str(e)}")
        
        if progress_callback:
            progress_callback((slide_idx + 1) / total_slides)
    
    # Summary of processing methods used
    log_messages.append(f"✓ Processed {images_processed} images ({vision_used_count} with vision, {text_fallback_count} text-only) and {tables_processed} tables")
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output, log_messages, alt_text_log


def create_zip_file(files_dict):
    """Create a ZIP file from a dictionary of {filename: bytes_io}"""
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for filename, file_bytes in files_dict.items():
            file_bytes.seek(0)
            zip_file.writestr(filename, file_bytes.read())
    zip_buffer.seek(0)
    return zip_buffer


def check_windows_platform():
    """Check if running on Windows platform"""
    import platform
    return platform.system() == "Windows"


def convert_pptx_to_pdf_subprocess(pptx_files_dict, temp_dir, generate_slides=False, generate_handouts=False, slides_per_page=3):
    """
    Convert PowerPoint files to PDF by calling external script via subprocess.
    This allows the COM automation to run in a separate process.
    """
    import subprocess
    import json
    
    slides_pdf_dict = {}
    handouts_pdf_dict = {}
    
    if not generate_slides and not generate_handouts:
        return slides_pdf_dict, handouts_pdf_dict
    
    # Map slides_per_page to PowerPoint output type constants
    # Slides = 1, 2 slides = 2, 3 slides = 3, 4 slides = 8, 6 slides = 4
    handout_output_map = {2: 2, 3: 3, 4: 8, 6: 4}
    handout_output_type = handout_output_map.get(slides_per_page, 3)
    
    for filename, pptx_bytes in pptx_files_dict.items():
        base_name = os.path.splitext(filename)[0]
        
        # Save PPTX to temp file
        temp_pptx_path = os.path.join(temp_dir, filename)
        pptx_bytes.seek(0)
        with open(temp_pptx_path, 'wb') as f:
            f.write(pptx_bytes.read())
        pptx_bytes.seek(0)
        
        try:
            # Call external PDF converter script
            result = subprocess.run(
                [
                    'python', 'pdf_converter.py',
                    temp_pptx_path,
                    temp_dir,
                    str(generate_slides).lower(),
                    str(generate_handouts).lower(),
                    str(slides_per_page)
                ],
                capture_output=True,
                text=True,
                timeout=120  # 2 minute timeout per file
            )
            
            # Parse JSON result
            if result.returncode == 0:
                conversion_result = json.loads(result.stdout)
                
                if conversion_result.get('success'):
                    # Read slides PDF if created
                    if conversion_result.get('slides_pdf') and os.path.exists(conversion_result['slides_pdf']):
                        with open(conversion_result['slides_pdf'], 'rb') as f:
                            slides_pdf_dict[f"{base_name}_slides.pdf"] = BytesIO(f.read())
                    
                    # Read handouts PDF if created
                    if conversion_result.get('handouts_pdf') and os.path.exists(conversion_result['handouts_pdf']):
                        with open(conversion_result['handouts_pdf'], 'rb') as f:
                            handouts_pdf_dict[f"{base_name}_handouts.pdf"] = BytesIO(f.read())
                else:
                    st.warning(f"PDF conversion failed for {filename}: {conversion_result.get('error', 'Unknown error')}")
            else:
                st.warning(f"PDF conversion process failed for {filename}: {result.stderr}")
                
        except subprocess.TimeoutExpired:
            st.warning(f"PDF conversion timed out for {filename}")
        except Exception as e:
            st.warning(f"Error converting {filename} to PDF: {str(e)}")
    
    return slides_pdf_dict, handouts_pdf_dict


# Main Streamlit App
st.title("PowerPoint Digital Accessibility Improvements and File Save Utility")

st.markdown("""
This application improves the digital accessibility of PowerPoint files by generating alt text
for images and tables and enables batch saving of PDF slides and handouts.
""")

api_key = st.text_input("TAMU API Key", value=None, type="password")

if api_key is not None and api_key != "":
    try:
        model_dict = call_models_api(api_key)
        selected_model_name = st.selectbox(
            "Pick a large language model to use for generating alt text", 
            list(model_dict), key=25
        )
        selected_model_id = model_dict[selected_model_name]
        
        # Selection for completed files
        st.subheader("Output Options")
        
        # Check if on Windows for PDF options
        is_windows = check_windows_platform()
        
        if is_windows:
            slides = st.checkbox("Select to output slides in PDF format")
            handouts = st.checkbox("Select to output handouts in PDF format")
            
            if handouts:
                slides_per_page = st.radio("Select number of slides per page", [2, 3, 4, 6], index=1)
            else:
                slides_per_page = 3
        else:
            st.info("PDF generation is only available on Windows with PowerPoint installed.")
            slides = False
            handouts = False
            slides_per_page = 3
                
        # File uploader
        uploaded_files = st.file_uploader(
            "Choose PowerPoint files to modify", 
            type="pptx", 
            accept_multiple_files=True
        )

        # Process button
        if uploaded_files:
            st.write(f"**{len(uploaded_files)} file(s) uploaded**")
            
            if st.button("Process PowerPoint Files", type="primary"):
                processed_pptx = {}
                processed_slides_pdf = {}
                processed_handouts_pdf = {}
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # Step 1: Process all PowerPoint files for accessibility
                with st.expander("Processing Details - Accessibility", expanded=True):
                    for i, uploaded_file in enumerate(uploaded_files):
                        status_text.text(f"Processing accessibility: {uploaded_file.name}")
                        st.write(f"**Processing: {uploaded_file.name}**")
                        
                        try:
                            # Process PowerPoint for accessibility
                            modified_pptx, log_messages, alt_text_log = process_powerpoint(
                                uploaded_file, api_key, selected_model_id, use_vision=True
                            )
                            
                            # Display log messages
                            for msg in log_messages:
                                st.write(msg)
                            
                            # Store processed PPTX
                            base_name = os.path.splitext(uploaded_file.name)[0]
                            processed_pptx[f"{base_name}_accessible.pptx"] = modified_pptx
                            
                            st.success(f"✓ Completed accessibility processing: {uploaded_file.name}")
                            
                        except Exception as e:
                            st.error(f"Error processing {uploaded_file.name}: {str(e)}")
                        
                        progress_bar.progress((i + 1) / len(uploaded_files) * 0.5)
                
                # Step 2: Generate PDFs if requested using subprocess
                if (slides or handouts) and is_windows:
                    with st.expander("Processing Details - PDF Generation", expanded=True):
                        status_text.text("Generating PDFs...")
                        st.write("**Generating PDF files using PowerPoint...**")
                        
                        # Create temp directory for PDF generation
                        temp_dir = tempfile.mkdtemp()
                        
                        try:
                            processed_slides_pdf, processed_handouts_pdf = convert_pptx_to_pdf_subprocess(
                                processed_pptx,
                                temp_dir,
                                generate_slides=slides,
                                generate_handouts=handouts,
                                slides_per_page=slides_per_page
                            )
                            
                            if slides:
                                st.write(f"✓ Generated {len(processed_slides_pdf)} slides PDF(s)")
                            if handouts:
                                st.write(f"✓ Generated {len(processed_handouts_pdf)} handouts PDF(s)")
                                
                        except Exception as e:
                            st.error(f"Error generating PDFs: {str(e)}")
                        
                        finally:
                            # Clean up temp files
                            try:
                                shutil.rmtree(temp_dir)
                            except:
                                pass
                        
                        progress_bar.progress(1.0)
                
                status_text.text("Processing complete!")
                progress_bar.progress(1.0)
                
                # Store in session state to persist after button click
                st.session_state['processed_pptx'] = processed_pptx
                st.session_state['processed_slides_pdf'] = processed_slides_pdf
                st.session_state['processed_handouts_pdf'] = processed_handouts_pdf
                st.session_state['processing_complete'] = True
        
        # Download section - check if processing is complete
        if st.session_state.get('processing_complete', False):
            st.subheader("📥 Download Processed Files")
            
            processed_pptx = st.session_state.get('processed_pptx', {})
            processed_slides_pdf = st.session_state.get('processed_slides_pdf', {})
            processed_handouts_pdf = st.session_state.get('processed_handouts_pdf', {})
            
            if len(processed_pptx) == 1:
                # Single file - provide individual downloads
                st.write("**Download individual files:**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    for filename, file_bytes in processed_pptx.items():
                        file_bytes.seek(0)
                        st.download_button(
                            label="📄 Download Accessible PPTX",
                            data=file_bytes,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="download_pptx"
                        )
                
                with col2:
                    if processed_slides_pdf:
                        for filename, file_bytes in processed_slides_pdf.items():
                            file_bytes.seek(0)
                            st.download_button(
                                label="📑 Download Slides PDF",
                                data=file_bytes,
                                file_name=filename,
                                mime="application/pdf",
                                key="download_slides"
                            )
                    else:
                        st.write("*No slides PDF*")
                
                with col3:
                    if processed_handouts_pdf:
                        for filename, file_bytes in processed_handouts_pdf.items():
                            file_bytes.seek(0)
                            st.download_button(
                                label="📋 Download Handouts PDF",
                                data=file_bytes,
                                file_name=filename,
                                mime="application/pdf",
                                key="download_handouts"
                            )
                    else:
                        st.write("*No handouts PDF*")
            
            elif len(processed_pptx) > 1:
                # Multiple files - provide ZIP downloads
                st.write("**Download ZIP archives:**")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if processed_pptx:
                        pptx_zip = create_zip_file(processed_pptx)
                        st.download_button(
                            label="📦 Download All PPTX (ZIP)",
                            data=pptx_zip,
                            file_name="accessible_presentations.zip",
                            mime="application/zip",
                            key="download_pptx_zip"
                        )
                
                with col2:
                    if processed_slides_pdf:
                        slides_zip = create_zip_file(processed_slides_pdf)
                        st.download_button(
                            label="📦 Download All Slides PDF (ZIP)",
                            data=slides_zip,
                            file_name="slides_pdfs.zip",
                            mime="application/zip",
                            key="download_slides_zip"
                        )
                    else:
                        st.write("*No slides PDFs*")
                
                with col3:
                    if processed_handouts_pdf:
                        handouts_zip = create_zip_file(processed_handouts_pdf)
                        st.download_button(
                            label="📦 Download All Handouts PDF (ZIP)",
                            data=handouts_zip,
                            file_name="handouts_pdfs.zip",
                            mime="application/zip",
                            key="download_handouts_zip"
                        )
                    else:
                        st.write("*No handouts PDFs*")
                
                # Option to download everything in one ZIP
                st.write("---")
                all_files = {**processed_pptx, **processed_slides_pdf, **processed_handouts_pdf}
                if all_files:
                    all_zip = create_zip_file(all_files)
                    st.download_button(
                        label="📦 Download All Files (Single ZIP)",
                        data=all_zip,
                        file_name="all_processed_files.zip",
                        mime="application/zip",
                        key="download_all_zip"
                    )
            
            # Reset button
            if st.button("🔄 Process New Files"):
                st.session_state['processing_complete'] = False
                st.session_state['processed_pptx'] = {}
                st.session_state['processed_slides_pdf'] = {}
                st.session_state['processed_handouts_pdf'] = {}
                st.rerun()
    except Exception as e:
        st.error(f"Error connecting to API: {str(e)}")
        st.write("Please check your API key and try again.")

else:
    st.warning("Please enter your TAMU API Key to continue.")
    st.markdown("""
    ### Instructions:
    1. Enter your TAMU AI API key above
    2. Select a language model for use in generating alt text
    3. Select additional formats for saving files
    4. Upload PowerPoint files to process
    5. Click 'Process PowerPoint Files' to make them accessible and save to multiple formats
    6. Download the modified files individually or as a ZIP
    """)            
